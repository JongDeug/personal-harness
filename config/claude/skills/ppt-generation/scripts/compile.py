#!/usr/bin/env python3
"""
compile.py — screenshots/slide*.png → <project>.pptx 컴파일

사용법:
  python3 compile.py                 # ./screenshots/slide*.png → ./presentation.pptx
  python3 compile.py my-deck.pptx    # 출력 파일명 지정
"""

import os
import sys
from pptx import Presentation
from pptx.util import Inches


def main():
    cwd = os.getcwd()
    screenshots_dir = os.path.join(cwd, "screenshots")
    output_path = sys.argv[1] if len(sys.argv) > 1 else os.path.join(cwd, "presentation.pptx")

    if not os.path.isdir(screenshots_dir):
        print(f"ERROR: screenshots dir not found: {screenshots_dir}")
        sys.exit(1)

    slide_images = sorted(
        os.path.join(screenshots_dir, f)
        for f in os.listdir(screenshots_dir)
        if f.startswith("slide") and f.endswith(".png")
    )

    if not slide_images:
        print(f"ERROR: no slide*.png in {screenshots_dir}")
        sys.exit(1)

    # 1920×1080 @ 96dpi == 20" × 11.25"
    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)
    blank_layout = prs.slide_layouts[6]

    for img in slide_images:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img, 0, 0, Inches(20), Inches(11.25))
        print(f"  added: {os.path.basename(img)}")

    prs.save(output_path)
    print(f"PPTX saved: {output_path}")


if __name__ == "__main__":
    main()
