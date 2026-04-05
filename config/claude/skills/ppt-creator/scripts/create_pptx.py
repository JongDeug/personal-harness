#!/usr/bin/env python3
"""
create_pptx.py - Generate a .pptx file from a JSON slide structure.

Usage:
  python3 create_pptx.py <input.json> [output.pptx] [--script]

Options:
  --script    Also generate a speaker script markdown file alongside the .pptx

JSON structure:
{
  "title": "Presentation Title",
  "subtitle": "Optional subtitle",
  "author": "Optional author",
  "theme": "dark" | "light" | "blue",   // optional, default "blue"
  "slides": [
    {
      "type": "title",
      "title": "...",
      "subtitle": "...",
      "notes": "발표자 노트 (선택)"
    },
    {
      "type": "section",
      "title": "...",
      "notes": "..."
    },
    {
      "type": "bullets",
      "title": "...",
      "bullets": ["point 1", "**bold point**", "- sub-point"],  // **text** for bold
      "notes": "..."
    },
    {
      "type": "code",
      "title": "...",
      "code": "...",
      "language": "python",
      "notes": "..."
    },
    {
      "type": "two_column",
      "title": "...",
      "left": { "heading": "...", "bullets": ["..."] },
      "right": { "heading": "...", "bullets": ["..."] },
      "notes": "..."
    },
    {
      "type": "table",
      "title": "...",
      "headers": ["Col1", "Col2", "Col3"],
      "rows": [["A", "B", "C"], ["D", "E", "F"]],
      "notes": "..."
    },
    {
      "type": "image",
      "title": "...",
      "image_path": "/path/to/image.png",
      "caption": "이미지 설명 (선택)",
      "notes": "..."
    },
    {
      "type": "quote",
      "quote": "...",
      "attribution": "...",
      "notes": "..."
    },
    {
      "type": "end",
      "title": "Thank you",
      "subtitle": "Q&A",
      "notes": "..."
    }
  ]
}
"""

import json
import sys
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Theme definitions ───────────────────────────────────────────────────────
THEMES = {
    "blue": {
        "bg_title":    RGBColor(0x1E, 0x3A, 0x5F),
        "bg_section":  RGBColor(0x2E, 0x5E, 0xA8),
        "bg_content":  RGBColor(0xFF, 0xFF, 0xFF),
        "bg_code":     RGBColor(0xF4, 0xF6, 0xF9),
        "accent":      RGBColor(0x2E, 0x5E, 0xA8),
        "text_light":  RGBColor(0xFF, 0xFF, 0xFF),
        "text_dark":   RGBColor(0x1A, 0x1A, 0x2E),
        "text_muted":  RGBColor(0x6B, 0x7A, 0x99),
        "code_text":   RGBColor(0x1A, 0x1A, 0x2E),
        "bullet_dot":  RGBColor(0x2E, 0x5E, 0xA8),
        "table_header_bg": RGBColor(0x2E, 0x5E, 0xA8),
        "table_row_even":  RGBColor(0xF0, 0xF4, 0xFF),
        "table_row_odd":   RGBColor(0xFF, 0xFF, 0xFF),
        "table_border":    RGBColor(0xCC, 0xD6, 0xE8),
        "is_dark_title":   True,
    },
    "dark": {
        "bg_title":    RGBColor(0x0D, 0x0D, 0x0D),
        "bg_section":  RGBColor(0x1A, 0x1A, 0x2E),
        "bg_content":  RGBColor(0x16, 0x16, 0x16),
        "bg_code":     RGBColor(0x0D, 0x0D, 0x0D),
        "accent":      RGBColor(0x00, 0xD4, 0xFF),
        "text_light":  RGBColor(0xFF, 0xFF, 0xFF),
        "text_dark":   RGBColor(0xE0, 0xE0, 0xE0),
        "text_muted":  RGBColor(0x88, 0x88, 0x88),
        "code_text":   RGBColor(0x00, 0xD4, 0xFF),
        "bullet_dot":  RGBColor(0x00, 0xD4, 0xFF),
        "table_header_bg": RGBColor(0x00, 0xD4, 0xFF),
        "table_row_even":  RGBColor(0x22, 0x22, 0x33),
        "table_row_odd":   RGBColor(0x1A, 0x1A, 0x2A),
        "table_border":    RGBColor(0x33, 0x33, 0x55),
        "is_dark_title":   True,
    },
    "light": {
        "bg_title":    RGBColor(0xF5, 0xF5, 0xF5),
        "bg_section":  RGBColor(0xE8, 0xE8, 0xE8),
        "bg_content":  RGBColor(0xFF, 0xFF, 0xFF),
        "bg_code":     RGBColor(0xF0, 0xF0, 0xF0),
        "accent":      RGBColor(0x33, 0x69, 0xE8),
        "text_light":  RGBColor(0x33, 0x33, 0x33),
        "text_dark":   RGBColor(0x1A, 0x1A, 0x1A),
        "text_muted":  RGBColor(0x77, 0x77, 0x77),
        "code_text":   RGBColor(0x1A, 0x1A, 0x1A),
        "bullet_dot":  RGBColor(0x33, 0x69, 0xE8),
        "table_header_bg": RGBColor(0x33, 0x69, 0xE8),
        "table_row_even":  RGBColor(0xF5, 0xF7, 0xFF),
        "table_row_odd":   RGBColor(0xFF, 0xFF, 0xFF),
        "table_border":    RGBColor(0xDD, 0xDD, 0xDD),
        "is_dark_title":   False,
    },
}

W = Inches(13.33)
H = Inches(7.5)
FONT_NAME = "Malgun Gothic"


# ── Helpers ─────────────────────────────────────────────────────────────────

def set_bg(slide, color: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, x, y, w, h, font_size=24, bold=False,
                color=RGBColor(0xFF, 0xFF, 0xFF), align=PP_ALIGN.LEFT,
                font_name=FONT_NAME, word_wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


def add_rect(slide, x, y, w, h, color: RGBColor):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def parse_inline_bold(text):
    """
    Parse **bold** markers in text.
    Returns list of (segment_text, is_bold) tuples.
    """
    parts = re.split(r'(\*\*[^*]+\*\*)', text)
    result = []
    for part in parts:
        if part.startswith('**') and part.endswith('**'):
            result.append((part[2:-2], True))
        elif part:
            result.append((part, False))
    return result


def add_rich_paragraph(tf, text, font_size, base_color, is_first=False,
                       space_before=None, is_sub=False):
    """
    Add a paragraph with inline bold support (**text**).
    """
    p = tf.paragraphs[0] if is_first else tf.add_paragraph()
    if space_before is not None:
        p.space_before = Pt(space_before)

    segments = parse_inline_bold(text)
    for seg_text, is_bold in segments:
        run = p.add_run()
        run.text = seg_text
        run.font.size = Pt(font_size)
        run.font.bold = is_bold
        run.font.color.rgb = base_color
        run.font.name = FONT_NAME
    return p


def set_notes(slide, notes_text):
    """Insert text into slide's speaker notes."""
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_slide_number(slide, num, t):
    """Add slide number to bottom-right corner."""
    is_dark = t.get("is_dark_title", False)
    # content slides use text_muted, dark bg slides use text_light
    color = t["text_muted"]
    add_textbox(slide, str(num),
                Inches(12.6), Inches(7.05), Inches(0.6), Inches(0.35),
                font_size=11, color=color, align=PP_ALIGN.RIGHT)


def add_title_bar(slide, title_text, t):
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.2), t["accent"])
    add_textbox(slide, title_text,
                Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.9),
                font_size=28, bold=True, color=t["text_light"],
                align=PP_ALIGN.LEFT)


# ── Slide builders ───────────────────────────────────────────────────────────

def build_title_slide(prs, data, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_title"])
    add_rect(slide, Inches(0), Inches(0), Inches(0.3), H, t["accent"])

    add_textbox(slide,
                data.get("title", "Untitled"),
                Inches(0.8), Inches(1.8), Inches(11.5), Inches(2.0),
                font_size=48, bold=True, color=t["text_light"],
                align=PP_ALIGN.LEFT)

    subtitle = data.get("subtitle", "")
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.8), Inches(3.9), Inches(11.5), Inches(1.0),
                    font_size=24, color=t["text_muted"],
                    align=PP_ALIGN.LEFT)

    author = data.get("author", "")
    if author:
        add_textbox(slide, author,
                    Inches(0.8), Inches(5.5), Inches(11.5), Inches(0.6),
                    font_size=16, color=t["text_muted"], align=PP_ALIGN.LEFT)

    set_notes(slide, data.get("notes", ""))
    return slide


def build_section_slide(prs, slide_data, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_section"])
    add_rect(slide, Inches(0), Inches(3.5), W, Inches(0.05), t["text_light"])

    add_textbox(slide,
                slide_data.get("title", "Section"),
                Inches(1.0), Inches(2.2), Inches(11.0), Inches(1.5),
                font_size=40, bold=True, color=t["text_light"],
                align=PP_ALIGN.CENTER)

    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(1.0), Inches(3.7), Inches(11.0), Inches(1.0),
                    font_size=20, color=t["text_light"],
                    align=PP_ALIGN.CENTER)

    set_notes(slide, slide_data.get("notes", ""))
    return slide


def build_bullets_slide(prs, slide_data, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_content"])
    add_title_bar(slide, slide_data.get("title", ""), t)

    bullets = slide_data.get("bullets", [])
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(12.0), Inches(5.6))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        is_sub = bullet.startswith("-")
        raw_text = bullet.lstrip("- ").strip()
        prefix = "    • " if is_sub else "• "
        text = prefix + raw_text
        fs = 18 if is_sub else 22
        sp = 4 if is_sub else 10

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(sp)

        segments = parse_inline_bold(text)
        for seg_text, is_bold in segments:
            run = p.add_run()
            run.text = seg_text
            run.font.size = Pt(fs)
            run.font.bold = is_bold
            run.font.color.rgb = t["text_dark"]
            run.font.name = FONT_NAME

    set_notes(slide, slide_data.get("notes", ""))
    return slide


def build_code_slide(prs, slide_data, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_content"])
    add_title_bar(slide, slide_data.get("title", ""), t)

    lang = slide_data.get("language", "")
    if lang:
        add_textbox(slide, lang.upper(),
                    Inches(11.4), Inches(0.2), Inches(1.5), Inches(0.8),
                    font_size=13, color=t["text_light"], align=PP_ALIGN.RIGHT)

    add_rect(slide, Inches(0.4), Inches(1.4), Inches(12.5), Inches(5.7), t["bg_code"])

    code = slide_data.get("code", "")
    txBox = slide.shapes.add_textbox(Inches(0.65), Inches(1.55), Inches(12.0), Inches(5.4))
    tf = txBox.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = code
    run.font.size = Pt(15)
    run.font.name = "Courier New"
    run.font.color.rgb = t["code_text"]

    set_notes(slide, slide_data.get("notes", ""))
    return slide


def build_two_column_slide(prs, slide_data, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_content"])
    add_title_bar(slide, slide_data.get("title", ""), t)

    add_rect(slide, Inches(6.6), Inches(1.4), Inches(0.05), Inches(5.7), t["accent"])

    for col, x in [("left", Inches(0.4)), ("right", Inches(6.9))]:
        col_data = slide_data.get(col, {})
        heading = col_data.get("heading", "")
        bullets = col_data.get("bullets", [])

        if heading:
            add_textbox(slide, heading, x, Inches(1.4), Inches(6.0), Inches(0.7),
                        font_size=20, bold=True, color=t["accent"])

        if bullets:
            txBox = slide.shapes.add_textbox(x, Inches(2.2), Inches(6.0), Inches(4.8))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(bullets):
                is_sub = bullet.startswith("-")
                raw_text = bullet.lstrip("- ").strip()
                prefix = "    • " if is_sub else "• "
                text = prefix + raw_text
                fs = 16 if is_sub else 18
                sp = 4 if is_sub else 8

                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(sp)

                segments = parse_inline_bold(text)
                for seg_text, is_bold in segments:
                    run = p.add_run()
                    run.text = seg_text
                    run.font.size = Pt(fs)
                    run.font.bold = is_bold
                    run.font.color.rgb = t["text_dark"]
                    run.font.name = FONT_NAME

    set_notes(slide, slide_data.get("notes", ""))
    return slide


def build_table_slide(prs, slide_data, t):
    """Table slide: headers + rows grid."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_content"])
    add_title_bar(slide, slide_data.get("title", ""), t)

    headers = slide_data.get("headers", [])
    rows = slide_data.get("rows", [])

    if not headers and not rows:
        set_notes(slide, slide_data.get("notes", ""))
        return slide

    num_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
    num_rows = len(rows) + (1 if headers else 0)

    table_x = Inches(0.5)
    table_y = Inches(1.45)
    table_w = Inches(12.3)
    table_h = Inches(5.6)

    table = slide.shapes.add_table(num_rows, num_cols, table_x, table_y, table_w, table_h).table

    col_w = Inches(12.3 / num_cols)
    for col_idx in range(num_cols):
        table.columns[col_idx].width = col_w

    def set_cell(cell, text, bg_color, text_color, bold=False, font_size=16):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = str(text)
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = text_color
        run.font.name = FONT_NAME

    row_offset = 0
    if headers:
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            set_cell(cell, header, t["table_header_bg"], t["text_light"],
                     bold=True, font_size=17)
        row_offset = 1

    for row_idx, row_data in enumerate(rows):
        bg = t["table_row_even"] if row_idx % 2 == 0 else t["table_row_odd"]
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + row_offset, col_idx)
            set_cell(cell, cell_text, bg, t["text_dark"], font_size=16)

    set_notes(slide, slide_data.get("notes", ""))
    return slide


def build_image_slide(prs, slide_data, t):
    """Image slide: title + embedded image + optional caption."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_content"])
    add_title_bar(slide, slide_data.get("title", ""), t)

    image_path = slide_data.get("image_path", "")
    caption = slide_data.get("caption", "")

    if image_path and Path(image_path).exists():
        img_top = Inches(1.4)
        img_h = Inches(5.4) if not caption else Inches(4.8)
        img_w = Inches(12.0)

        # Center image, preserve aspect ratio
        from PIL import Image as PILImage
        try:
            with PILImage.open(image_path) as img:
                orig_w, orig_h = img.size
            ratio = orig_w / orig_h
            calc_w = img_h * ratio
            if calc_w > img_w:
                calc_w = img_w
                img_h = calc_w / ratio
            left = Inches(0.5) + (img_w - calc_w) / 2
            slide.shapes.add_picture(image_path, left, img_top, calc_w, img_h)
        except Exception:
            # Pillow not available or error — add as-is
            slide.shapes.add_picture(image_path,
                                     Inches(0.65), img_top,
                                     Inches(12.0), Inches(4.8))

        if caption:
            add_textbox(slide, caption,
                        Inches(0.6), Inches(6.4), Inches(12.1), Inches(0.7),
                        font_size=14, color=t["text_muted"],
                        align=PP_ALIGN.CENTER)
    else:
        # Placeholder when image not found
        add_rect(slide, Inches(0.65), Inches(1.5), Inches(12.0), Inches(4.8),
                 t["bg_code"])
        add_textbox(slide,
                    f"[이미지 없음: {image_path}]",
                    Inches(0.65), Inches(3.5), Inches(12.0), Inches(0.8),
                    font_size=16, color=t["text_muted"],
                    align=PP_ALIGN.CENTER)

    set_notes(slide, slide_data.get("notes", ""))
    return slide


def build_quote_slide(prs, slide_data, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_title"])

    add_textbox(slide, "\u201c",
                Inches(0.5), Inches(0.3), Inches(2.0), Inches(2.0),
                font_size=120, color=t["accent"], align=PP_ALIGN.LEFT)

    add_textbox(slide,
                slide_data.get("quote", ""),
                Inches(1.2), Inches(1.5), Inches(11.0), Inches(3.5),
                font_size=28, bold=True, color=t["text_light"],
                align=PP_ALIGN.LEFT)

    attr = slide_data.get("attribution", "")
    if attr:
        add_textbox(slide, f"— {attr}",
                    Inches(1.2), Inches(5.2), Inches(11.0), Inches(0.8),
                    font_size=18, color=t["text_muted"],
                    align=PP_ALIGN.RIGHT)

    set_notes(slide, slide_data.get("notes", ""))
    return slide


def build_end_slide(prs, slide_data, t):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, t["bg_title"])
    add_rect(slide, Inches(0), Inches(0), Inches(0.3), H, t["accent"])

    add_textbox(slide,
                slide_data.get("title", "Thank You"),
                Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.5),
                font_size=52, bold=True, color=t["text_light"],
                align=PP_ALIGN.LEFT)

    subtitle = slide_data.get("subtitle", "")
    if subtitle:
        add_textbox(slide, subtitle,
                    Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.0),
                    font_size=24, color=t["text_muted"],
                    align=PP_ALIGN.LEFT)

    set_notes(slide, slide_data.get("notes", ""))
    return slide


# ── Speaker script generator ─────────────────────────────────────────────────

def generate_script(data, output_path):
    """Generate a markdown speaker script from notes fields."""
    lines = []
    title = data.get("title", "Presentation")
    lines.append(f"# 발표 대본: {title}\n")
    lines.append(f"> 자동 생성된 발표자 스크립트\n")
    lines.append("---\n")

    slides = data.get("slides", [])
    script_count = 0

    for i, slide_data in enumerate(slides, 1):
        slide_title = slide_data.get("title", slide_data.get("quote", f"슬라이드 {i}"))
        slide_type = slide_data.get("type", "bullets")
        notes = slide_data.get("notes", "")

        lines.append(f"## [{i}] {slide_title}")
        lines.append(f"*타입: {slide_type}*\n")

        if notes:
            lines.append(notes)
            script_count += 1
        else:
            lines.append("*(발표자 노트 없음)*")

        lines.append("\n---\n")

    script_path = str(output_path).replace(".pptx", "_script.md")
    with open(script_path, "w", encoding="utf-8") as f:
        f.write("\n".join(lines))

    return script_path, script_count


# ── Main ─────────────────────────────────────────────────────────────────────

SLIDE_BUILDERS = {
    "title":      build_title_slide,
    "section":    build_section_slide,
    "bullets":    build_bullets_slide,
    "code":       build_code_slide,
    "two_column": build_two_column_slide,
    "table":      build_table_slide,
    "image":      build_image_slide,
    "quote":      build_quote_slide,
    "end":        build_end_slide,
}

# Slides that use dark backgrounds (slide numbers shown lighter)
DARK_BG_TYPES = {"title", "section", "quote", "end"}


def create_pptx(input_path: str, output_path: str, gen_script: bool = False):
    with open(input_path, "r", encoding="utf-8") as f:
        data = json.load(f)

    theme_name = data.get("theme", "blue")
    t = THEMES.get(theme_name, THEMES["blue"])

    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slides = data.get("slides", [])
    built_count = 0
    slide_num_counter = 0

    for slide_data in slides:
        slide_type = slide_data.get("type", "bullets")
        builder = SLIDE_BUILDERS.get(slide_type)

        if builder is None:
            print(f"[WARN] Unknown slide type: '{slide_type}', skipping.")
            continue

        slide = builder(prs, slide_data, t)
        built_count += 1
        slide_num_counter += 1

        # Slide numbers: skip title and end slides
        if slide_type not in ("title", "end"):
            add_slide_number(slide, slide_num_counter, t)

    prs.save(output_path)
    print(f"[OK] Saved: {output_path} ({built_count} slides)")

    if gen_script:
        script_path, notes_count = generate_script(data, Path(output_path))
        print(f"[OK] Script: {script_path} ({notes_count} slides with notes)")

    return built_count


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 create_pptx.py <input.json> [output.pptx] [--script]")
        sys.exit(1)

    args = sys.argv[1:]
    gen_script_flag = "--script" in args
    args = [a for a in args if a != "--script"]

    input_file = args[0]
    output_file = args[1] if len(args) > 1 else Path(args[0]).stem + ".pptx"
    create_pptx(input_file, output_file, gen_script=gen_script_flag)
